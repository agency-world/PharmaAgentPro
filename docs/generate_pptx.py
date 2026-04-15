"""Generate MARKETING_DECK.pptx from the marketing deck content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path


# ── Brand Colors ──
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
PRIMARY = RGBColor(0x3B, 0x82, 0xF6)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xE2, 0xE8, 0xF0)
DIM_TEXT = RGBColor(0x94, 0xA3, 0xB8)
NAVY = RGBColor(0x1A, 0x36, 0x5D)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=LIGHT_TEXT, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_table_shape(slide, left, top, width, height, headers, rows,
                    header_color=PRIMARY, row_color=SURFACE, text_color=LIGHT_TEXT):
    num_rows = 1 + len(rows)
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols,
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(height))
    table = table_shape.table

    # Style headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    # Style rows
    for r_idx, row in enumerate(rows):
        fill_color = row_color if r_idx % 2 == 0 else DARK_BG
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = text_color
                p.font.name = "Calibri"

    return table_shape


def add_card(slide, left, top, width, height, title, body, accent_color=ACCENT):
    """Add a rounded rectangle card with title and body."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(10)
    p2.font.color.rgb = DIM_TEXT
    p2.font.name = "Calibri"
    p2.space_before = Pt(4)

    return shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Logo box
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5))
    logo.fill.solid()
    logo.fill.fore_color.rgb = PRIMARY
    logo.line.fill.background()
    tf = logo.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "P"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(10)

    add_textbox(slide, 2.5, 3.3, 8.3, 1, "PharmaAgent Pro", 48, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 2.5, 4.2, 8.3, 0.6, "Accelerating Drug Development with Agentic AI", 22, ACCENT, False, PP_ALIGN.CENTER)
    add_textbox(slide, 2.5, 5.0, 8.3, 0.5, "Built on Claude's Managed Agent Platform", 16, DIM_TEXT, False, PP_ALIGN.CENTER)
    add_textbox(slide, 2.5, 6.2, 8.3, 0.4, "Confidential  |  April 2026", 12, DIM_TEXT, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # SLIDE 2: THE PROBLEM
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "The Problem: Drug Development Is Broken", 32, WHITE, True)
    add_textbox(slide, 0.8, 1.1, 11.7, 0.5, "$2.6 Billion & 12 Years -- the average cost and time to bring a single drug to market.", 16, WARNING, False)

    # Pipeline visual
    pipeline_stages = [
        ("Compound\nDiscovery", "3-6 yr", "10,000\ncompounds"),
        ("Preclinical\nTesting", "1-2 yr", "250 enter\npreclinical"),
        ("Phase I\nTrial", "1-3 yr", "5 enter\nclinical"),
        ("Phase II\nTrial", "2-3 yr", "---"),
        ("Phase III\nTrial", "1-2 yr", "1 enters\nPhase III"),
        ("FDA Review\n& Approval", "1-2 yr", "1\nAPPROVED"),
    ]
    x_start = 0.8
    card_w = 1.85
    for i, (stage, duration, count) in enumerate(pipeline_stages):
        x = x_start + i * 2.05
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(1.9), Inches(card_w), Inches(1.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = PRIMARY if i < 5 else SUCCESS
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = stage
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY if i < 5 else SUCCESS
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = duration
        p2.font.size = Pt(10)
        p2.font.color.rgb = DIM_TEXT
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = count
        p3.font.size = Pt(9)
        p3.font.color.rgb = LIGHT_TEXT
        p3.alignment = PP_ALIGN.CENTER

    # Problem cards
    problems = [
        ("Data Silos", "Compound data, trial results, and safety signals live in disconnected systems"),
        ("Manual Documents", "Regulatory teams spend 40%+ of time on formatting, not strategy"),
        ("Delayed Safety", "Adverse event patterns spotted weeks late due to manual review"),
        ("Knowledge Fragmentation", "Critical drug interaction data scattered across teams"),
        ("Compliance Burden", "FDA/ICH requirements are error-prone when manual"),
    ]
    for i, (title, desc) in enumerate(problems):
        x = 0.8 + i * 2.45
        add_card(slide, x, 3.7, 2.25, 1.3, title, desc, DANGER)

    # Cost callout
    add_textbox(slide, 0.8, 5.3, 11.7, 0.4, "The Cost of Delay", 18, WHITE, True)
    cost_items = [
        "Every month of delay = $30-100M in lost revenue",
        "Late-stage trial failures from missed signals = $800M-1.5B wasted",
        "FDA rejection from incomplete submissions = 6-12 months rework",
    ]
    add_bullet_list(slide, 0.8, 5.8, 11.7, 1.5, cost_items, 14, LIGHT_TEXT)

    # ════════════════════════════════════════════
    # SLIDE 3: THE SOLUTION
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "The Solution: PharmaAgent Pro", 32, WHITE, True)
    add_textbox(slide, 0.8, 1.0, 11.7, 0.5, "An AI agent that thinks like your best researcher", 18, ACCENT, False)

    # Workflow demo
    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8), Inches(1.7), Inches(7.5), Inches(2.5))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = SURFACE
    demo_box.line.color.rgb = ACCENT
    demo_box.line.width = Pt(1)
    tf = demo_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    lines = [
        ('User: "Analyze Oncolytin-B safety and generate a CIOMS report"', LIGHT_TEXT, True),
        ("", WHITE, False),
        ("1. Routes to Safety Agent (Opus 4.6 + Extended Thinking)", ACCENT, False),
        ("2. Retrieves adverse event data via safety_signal_analysis tool", ACCENT, False),
        ("3. Detects ILD signal -- ALERT status, DSMB review recommended", WARNING, False),
        ("4. Generates ICH E2A-compliant CIOMS safety report", ACCENT, False),
        ("5. Validates against compliance checklist", SUCCESS, False),
        ("", WHITE, False),
        ("Complete safety analysis + CIOMS report in < 60 seconds", SUCCESS, True),
        ("(Previously: 2-3 days of manual work)", DIM_TEXT, False),
    ]
    for i, (text, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"

    # 4 Agent cards
    agents = [
        ("Drug Discovery Agent", "Compound profiling,\ndrug-likeness, CYP\ninteractions, SAR", "10x faster screening"),
        ("Clinical Trials Agent", "Trial status, enrollment,\nendpoint analysis,\ninterim results", "Real-time intelligence"),
        ("Regulatory Agent", "FDA/ICH-compliant\ndocument generation\n(IND, CSR, CIOMS)", "80% less prep time"),
        ("Safety Agent", "AE analysis, signal\ndetection, benefit-risk\nassessment", "Seconds, not weeks"),
    ]
    for i, (title, desc, metric) in enumerate(agents):
        x = 0.8 + i * 3.1
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(4.5), Inches(2.9), Inches(2.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = PRIMARY
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT_TEXT
        p2.font.name = "Calibri"
        p2.space_before = Pt(8)
        p3 = tf.add_paragraph()
        p3.text = metric
        p3.font.size = Pt(12)
        p3.font.bold = True
        p3.font.color.rgb = SUCCESS
        p3.font.name = "Calibri"
        p3.space_before = Pt(12)

    # ════════════════════════════════════════════
    # SLIDE 4: HOW IT WORKS
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "How It Works: Intelligent Query Routing", 32, WHITE, True)

    # Flow diagram with shapes
    flow_items = [
        ("User Question", LIGHT_TEXT, 1.0),
        ("AI Classifier\n(Haiku 4.5, <100ms)", WARNING, 2.0),
        ("Specialized Agent\n(Opus 4.6 + Thinking)", PRIMARY, 3.2),
        ("Tool Execution\n(Local Data Access)", ACCENT, 4.4),
        ("Synthesized Response\n(Compliant + Audited)", SUCCESS, 5.6),
    ]
    for label, color, y in flow_items:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.8), Inches(y), Inches(4.5), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    # Arrows between
    for y in [1.95, 3.15, 4.35]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                       Inches(2.8), Inches(y), Inches(0.3), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DIM_TEXT
        arrow.line.fill.background()

    # Key technology table on the right
    add_textbox(slide, 6.0, 1.0, 6.5, 0.5, "Key Technology Decisions", 20, WHITE, True)
    add_table_shape(slide, 6.0, 1.6, 6.5, 3.0,
        ["Feature", "How It Works", "Why It Matters"],
        [
            ["Multi-Model AI", "Haiku routes to Opus", "95% routing cost savings"],
            ["Extended Thinking", "Claude shows reasoning chain", "Transparent, auditable AI"],
            ["Tool-Based Retrieval", "AI selects tools, fetches data", "Precise, not black-box RAG"],
            ["Prompt Caching", "Reuse cached context", "90% cost reduction on follow-ups"],
            ["Guardrails", "PII redaction, injection defense", "Pharma-grade safety"],
        ]
    )

    # ════════════════════════════════════════════
    # SLIDE 5: PRODUCT SCREENSHOTS
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Product Experience", 32, WHITE, True)

    # Terminal screenshot mockup
    terminal = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8))
    terminal.fill.solid()
    terminal.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    terminal.line.color.rgb = DIM_TEXT
    terminal.line.width = Pt(1)
    tf = terminal.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)

    term_lines = [
        ("Terminal Chat Interface", ACCENT, True, 14),
        ("", WHITE, False, 10),
        ("PharmaAgent Pro -- Managed Agent Chatbot", DIM_TEXT, False, 10),
        ("Session: sesn_011Ca2Lmvdp...  Agent: agent_011Ca1Foj3X...", DIM_TEXT, False, 9),
        ("", WHITE, False, 10),
        ("You > Which compounds target EGFR?", SUCCESS, True, 11),
        ("", WHITE, False, 10),
        ("  (thinking...)", WARNING, False, 10),
        ("  Custom Tool: drug_lookup", ACCENT, False, 10),
        ("", WHITE, False, 10),
        ("Based on our compound library, Oncolytin-B (PHA-002)", LIGHT_TEXT, False, 10),
        ("targets EGFR:", LIGHT_TEXT, False, 10),
        ("", WHITE, False, 10),
        ("  - Compound: Oncolytin-B (PHA-002)", LIGHT_TEXT, False, 10),
        ("  - Class: Kinase Inhibitor", LIGHT_TEXT, False, 10),
        ("  - Target: EGFR/HER2 dual inhibitor", LIGHT_TEXT, False, 10),
        ("  - Phase: Phase III (ONCOLYZE-3 trial)", LIGHT_TEXT, False, 10),
        ("  - IC50: 2.3 nM (EGFR), 8.1 nM (HER2)", LIGHT_TEXT, False, 10),
        ("  - Indication: Non-Small Cell Lung Cancer", LIGHT_TEXT, False, 10),
    ]
    for i, (text, color, bold, size) in enumerate(term_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Courier New"

    # Web UI mockup
    web = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.8))
    web.fill.solid()
    web.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    web.line.color.rgb = DIM_TEXT
    web.line.width = Pt(1)
    tf = web.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)

    web_lines = [
        ("Web Chat Interface", ACCENT, True, 14),
        ("", WHITE, False, 10),
        ("[P] PharmaAgent Pro          * Connected", DIM_TEXT, False, 10),
        ("", WHITE, False, 10),
        ("YOU", PRIMARY, True, 10),
        ("Run safety signal analysis on Oncolytin-B", LIGHT_TEXT, False, 10),
        ("", WHITE, False, 10),
        ("PHARMAAGENT PRO", ACCENT, True, 10),
        ("[safety_signal_analysis] [adverse_event_search]", SUCCESS, False, 9),
        ("", WHITE, False, 10),
        ("Safety Signal Analysis: Oncolytin-B", WHITE, True, 11),
        ("", WHITE, False, 10),
        ("Overall Signal Status: ALERT", DANGER, True, 11),
        ("", WHITE, False, 10),
        ("  Total AEs:    3", LIGHT_TEXT, False, 10),
        ("  Serious AEs:  2 (ILD, QTc Prolongation)", LIGHT_TEXT, False, 10),
        ("  Grade 3+:     2", LIGHT_TEXT, False, 10),
        ("  Signal:       ILD (Respiratory SOC)", WARNING, False, 10),
        ("", WHITE, False, 10),
        ("Recommendation: Immediate DSMB review", DANGER, True, 10),
    ]
    for i, (text, color, bold, size) in enumerate(web_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Courier New"

    # ════════════════════════════════════════════
    # SLIDE 6: BENEFITS & IMPACT
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Benefits & Potential Impact", 32, WHITE, True)

    # Before/After metrics
    metrics = [
        ("Safety Signal Detection", "2-4 weeks", "< 60 seconds", "99%"),
        ("CIOMS Report Generation", "2-3 days", "< 2 minutes", "99%"),
        ("Compound Interaction Check", "4-6 hours", "< 30 seconds", "99%"),
        ("IND Annual Report Draft", "3-5 days", "10-15 minutes", "97%"),
    ]

    add_textbox(slide, 0.8, 1.1, 5.5, 0.4, "Quantified Value", 20, WHITE, True)
    add_table_shape(slide, 0.8, 1.6, 6.0, 2.5,
        ["Task", "Before", "After", "Faster"],
        [[m[0], m[1], m[2], m[3] + " faster"] for m in metrics]
    )

    # Stakeholder impact
    add_textbox(slide, 7.3, 1.1, 5.5, 0.4, "Impact by Stakeholder", 20, WHITE, True)
    stakeholders = [
        ("Medicinal Chemists", "10x faster lead screening"),
        ("Clinical Operations", "Eliminate weekly report cycles"),
        ("Regulatory Affairs", "80% less document prep time"),
        ("Pharmacovigilance", "Same-day signal detection"),
        ("Medical Writing", "Focus on content, not formatting"),
        ("Executive Leadership", "Data-driven pipeline decisions"),
    ]
    add_table_shape(slide, 7.3, 1.6, 5.5, 2.5,
        ["Stakeholder", "Key Benefit"],
        stakeholders
    )

    # Safety by design section
    add_textbox(slide, 0.8, 4.4, 11.7, 0.4, "Safety & Compliance by Design", 20, WHITE, True)

    safety_cards = [
        ("Zero PII Exposure", "Patient data detected and\nredacted before any API call.\nNever leaves the client.", SUCCESS),
        ("Full Audit Trail", "Every query, tool call, and\nresponse logged to immutable\nJSONL audit trail.", PRIMARY),
        ("FDA/ICH Validation", "Documents validated against\n21 CFR 312, ICH E2A/E2F/E3\nbefore delivery.", ACCENT),
        ("Rate Limiting", "Per-user sliding window\nprevents abuse and\ncontrols costs.", WARNING),
    ]
    for i, (title, desc, color) in enumerate(safety_cards):
        x = 0.8 + i * 3.1
        add_card(slide, x, 4.9, 2.9, 1.8, title, desc, color)

    # ════════════════════════════════════════════
    # SLIDE 7: COMPETITIVE ADVANTAGE
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Why PharmaAgent Pro?", 32, WHITE, True)

    add_textbox(slide, 0.8, 1.1, 5.8, 0.4, "vs. General-Purpose AI Chatbots", 18, WHITE, True)
    add_table_shape(slide, 0.8, 1.6, 5.8, 3.0,
        ["Capability", "Generic AI", "PharmaAgent Pro"],
        [
            ["Pharma-specific tools", "No", "6 specialized tools"],
            ["Regulatory compliance", "No", "Built-in ICH/FDA validation"],
            ["PII protection", "Basic", "Multi-layer detect + redact"],
            ["Audit trail", "No", "Immutable JSONL logging"],
            ["Extended thinking", "No", "Transparent reasoning"],
            ["Multi-agent routing", "No", "4 specialized agents"],
            ["Cost tracking", "No", "Per-query token + cost"],
        ]
    )

    add_textbox(slide, 7.0, 1.1, 5.8, 0.4, "vs. Traditional Pharma Software", 18, WHITE, True)
    add_table_shape(slide, 7.0, 1.6, 5.8, 3.0,
        ["Capability", "Legacy Tools", "PharmaAgent Pro"],
        [
            ["Natural language queries", "No", "Full conversational"],
            ["Cross-domain analysis", "Siloed", "Unified, 4 domains"],
            ["Setup time", "Months", "Minutes (1 API key)"],
            ["Document generation", "Templates only", "AI-generated + context"],
            ["Learning curve", "High", "Plain English queries"],
        ]
    )

    # ════════════════════════════════════════════
    # SLIDE 8: CASE STUDY
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Case Study: Safety Signal to Regulatory Report", 28, WHITE, True)
    add_textbox(slide, 0.8, 1.0, 11.7, 0.4, "Oncolytin-B ILD Signal Detection -- 3 minutes vs. 3-5 days", 16, ACCENT, False)

    # Three step cards
    steps = [
        ("Step 1: Signal Detection",
         "30 seconds",
         "User asks for safety analysis\n\nPharmaAgent Pro:\n- Calls safety_signal_analysis\n- Analyzes 3 AEs, identifies 2 SAEs\n- Flags ILD (Respiratory SOC)\n- Status: ALERT",
         WARNING),
        ("Step 2: Deep Dive",
         "45 seconds",
         "User requests SAE narratives\n\nPharmaAgent Pro:\n- Calls adverse_event_search\n- Returns 2 SAE records\n- ILD: Grade 3, drug discontinued\n- QTc: Grade 2, dose reduced",
         PRIMARY),
        ("Step 3: Regulatory Report",
         "90 seconds",
         "User requests CIOMS report\n\nPharmaAgent Pro:\n- Calls generate_document\n- ICH E2A-compliant CIOMS I\n- Patient, AE, drug, narrative\n- Validated compliance checklist",
         SUCCESS),
    ]
    for i, (title, time_text, desc, color) in enumerate(steps):
        x = 0.8 + i * 4.1
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(1.6), Inches(3.8), Inches(5.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.text = time_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.font.name = "Calibri"
        p2.space_before = Pt(8)
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = LIGHT_TEXT
        p3.font.name = "Calibri"
        p3.space_before = Pt(12)

    # Total time callout
    add_textbox(slide, 3.5, 6.8, 6.3, 0.5, "Total: ~3 minutes  |  Traditional: 3-5 days", 18, SUCCESS, True, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════
    # SLIDE 9: FEATURE ROADMAP
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.4, 11.7, 0.6, "Feature Roadmap", 32, WHITE, True)

    phases = [
        ("NOW\nPhase 1: MVP",
         ["6 pharma tools", "4 domain agents", "Guardrails & compliance", "Audit trail & tracking", "Managed Agent deploy", "Web + CLI chatbot"],
         SUCCESS, "LIVE"),
        ("Q3 2026\nPhase 2: Multi-Agent",
         ["Parallel agent execution", "Agent-to-agent comms", "Supervisor agent", "Dynamic tool discovery", "Conflict resolution"],
         PRIMARY, "PLANNED"),
        ("Q4 2026\nPhase 3: Enterprise",
         ["RBAC & multi-tenant", "FDA FAERS API live", "Document versioning", "Approval workflows", "Scheduled monitoring"],
         ACCENT, "PLANNED"),
        ("2027\nPhase 4: Advanced AI",
         ["Literature RAG", "Predictive safety", "Auto IND/NDA assembly", "Multi-modal analysis", "Federated learning"],
         WARNING, "FUTURE"),
    ]
    for i, (title, items, color, status) in enumerate(phases):
        x = 0.5 + i * 3.2
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(1.2), Inches(3.0), Inches(5.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.15)

        # Title
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Calibri"

        # Status badge
        p_status = tf.add_paragraph()
        p_status.text = status
        p_status.font.size = Pt(10)
        p_status.font.bold = True
        p_status.font.color.rgb = SUCCESS if status == "LIVE" else DIM_TEXT
        p_status.font.name = "Calibri"
        p_status.space_before = Pt(8)

        # Items
        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = f"  {item}"
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = LIGHT_TEXT
            p_item.font.name = "Calibri"
            p_item.space_before = Pt(6)

    # ════════════════════════════════════════════
    # SLIDE 10: CLOSING
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, 1.5, 1.5, 10.3, 0.8, "PharmaAgent Pro", 44, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 2.5, 10.3, 0.6, "From Compound to Clinic -- Faster, Safer, Smarter", 22, ACCENT, False, PP_ALIGN.CENTER)

    # Three value pillars
    pillars = [
        ("FASTER", "Signal detection in\nseconds, not weeks.\nDocuments in minutes,\nnot days.", SUCCESS),
        ("SAFER", "Zero PII exposure.\nFull audit trail.\nFDA/ICH compliant\nby design.", PRIMARY),
        ("SMARTER", "AI reasons through\ncomplex drug interactions.\nMulti-agent specialization\nfor every domain.", ACCENT),
    ]
    for i, (title, desc, color) in enumerate(pillars):
        x = 1.5 + i * 3.7
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(3.5), Inches(3.3), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = LIGHT_TEXT
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    add_textbox(slide, 1.5, 6.3, 10.3, 0.5, "Built on Claude  |  Engineered for Pharma  |  Ready for Production", 16, DIM_TEXT, False, PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    print("Generating MARKETING_DECK.pptx...")
    prs = build_presentation()
    output = Path(__file__).parent / "PharmaAgent_Pro_Marketing_Deck.pptx"
    prs.save(str(output))
    print(f"Saved to {output}")
